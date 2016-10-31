#import libraries to read pdfs
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from cStringIO import StringIO

#import library to read PPTX files
from pptx import Presentation

import os
import re

import psycopg2

def convert_pdf_to_txt(path):
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
    fp = file(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    maxpages = 0
    caching = True
    pagenos=set()

    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password, caching=caching, check_extractable=True):
        interpreter.process_page(page)

    text = retstr.getvalue()

    fp.close()
    device.close()
    retstr.close()
    return text

def pdfstringextractor( pdffilename ):
   # Choose PDF
   pdf2string = convert_pdf_to_txt(str(pdffilename))
   #pdf2string = re.sub("\n", "", pdf2string)
   pdf2string = pdf2string.replace("\n", "")
   return pdf2string;

def pptstringextractor( pptfilename ):
   prs = Presentation(str(pptfilename))

   # text_runs will be populated with a list of strings,
   # one for each text run in presentation
   ppt2string = []

   for slide in prs.slides:
       for shape in slide.shapes:
           if not shape.has_text_frame:
               continue
           for paragraph in shape.text_frame.paragraphs:
               for run in paragraph.runs:
                   ppt2string.append(run.text)

   return str(ppt2string);

os.chdir('C:\\Users\<<user name>>\Documents\DW') #Enter the complete path of the data warehouse folder
directory = os.getcwd()

conn = psycopg2.connect("dbname=<<enter postgresql database name>> host=localhost user=<<enter postgresql username>> password=<<enter postgresql password>>")
cur = conn.cursor()

# Extract PDFs		
for filename in os.listdir(directory):
    if filename.endswith(".pdf"): 
        pdfcheck = "true"
        pptcheck = "false"

        completepdftext = pdfstringextractor(filename)

        cur.execute("INSERT INTO test (dwfilename, ispdf, isppt, completetext) VALUES (%s, %s, %s, %s)",(str(filename), pdfcheck, pptcheck, completepdftext))
        conn.commit()

        continue
    else:
        continue

# Extract PPTs		
for filename in os.listdir(directory):
    if filename.endswith(".pptx"):
        pdfcheck = "false"
        pptcheck = "true"

        completeppttext = pptstringextractor(filename)

        cur.execute("INSERT INTO test (dwfilename, ispdf, isppt, completetext) VALUES (%s, %s, %s, %s)",(str(filename), pdfcheck, pptcheck, completeppttext))
        conn.commit()

        continue
    else:
        continue