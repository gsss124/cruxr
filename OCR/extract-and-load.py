from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from cStringIO import StringIO

from pptx import Presentation

import os
import re

import psycopg2

import sys
from PIL import Image
import pytesseract

import OleFileIO_PL as OleFile
import zipfile
import struct

startmark = "\xff\xd8"
startfix = 0
endmark = "\xff\xd9"
endfix = 2

DEBUG = False
CWD = '.'
CHUNK = 1024 * 64

# MS-ODRAW spec
formats = {
    # 2.2.24
    (0xF01A, 0x3D40): (50, ".emf"),
    (0xF01A, 0x3D50): (66, ".emf"),
    # 2.2.25
    (0xF01B, 0x2160): (50, ".wmf"),
    (0xF01B, 0x2170): (66, ".wmf"),
    # 2.2.26
    (0xF01C, 0x5420): (50, ".pict"),
    (0xF01C, 0x5430): (50, ".pict"),
    # 2.2.27
    (0xF01D, 0x46A0): (17, ".jpeg"),
    (0xF01D, 0x6E20): (17, ".jpeg"),
    (0xF01D, 0x46B0): (33, ".jpeg"),
    (0xF01D, 0x6E30): (33, ".jpeg"),
    # 2.2.28
    (0xF01E, 0x6E00): (17, ".png"),
    (0xF01E, 0x6E10): (33, ".png"),
    # 2.2.29
    (0xF01F, 0x7A80): (17, ".dib"),
    (0xF01F, 0x7A90): (33, ".dib"),
    # 2.2.30
    (0xF029, 0x6E40): (17, ".tiff"),
    (0xF029, 0x6E50): (33, ".tiff")
}

class InvalidFormat(Exception):
    pass


class PowerPointFormat(object):
    def __init__(self, filename):
        """
        filename:   archivo a abrir
        """
        self._files = {}
        
        # nombre base de imagenes
        self.basename = os.path.splitext(os.path.basename(filename))[0]
        
        self._process(filename)

    def extract(self, name, newname="", path=CWD):
        """
        Extrae imagen en directorio especificado.
        """
        self._extract(name, newname=newname, path=path)
        
    def extractall(self, path=CWD):
        """
        Extrae todas las imagenes en directorio especificado.
        """
        for img in self._files:
            self.extract(img, path=path)
            
    def namelist(self):
        """
        Retorna lista de imagenes contenidas en el archivo.
        """
        return self._files.keys()
        
    def __len__(self):
        return len(self._files)
        
    def __str__(self):
        return "<PowerPoint file with %s images>" % len(self)
    
    __repr__ = __str__


class PPT(PowerPointFormat):

    headerlen = struct.calcsize('<HHL')

    @classmethod
    def is_valid_format(cls, filename):
        return OleFile.isOleFile(filename)
    
    def _process(self, filename):

        olefile = OleFile.OleFileIO(filename)
        
        if not olefile.exists("Pictures"):
            return
        
        self.__stream = olefile.openstream("Pictures")

        stream = self.__stream
        offset = 0
        n = 1

        while True:
            header = stream.read(self.headerlen)
            offset += self.headerlen

            if not header: break
            
            recInstance, recType, recLen = struct.unpack_from("<HHL", header)

            stream.seek(recLen, 1)

            if DEBUG:
                print "%X %X %sb" % (recType, recInstance, recLen)
            
            extrabytes, ext = formats.get((recType, recInstance))
            
            recLen -= extrabytes
            offset += extrabytes
            
            filename = "{0}{1}{2}".format(self.basename, n, ext)
            
            self._files[filename] = (offset, recLen)
            offset += recLen
            
            n += 1

    def _extract(self, name, newname="", path=CWD):

        filename = newname or name
        
        if not name in self._files:
            raise IOError("No such file")
        
        offset, size = self._files[name]
        
        filepath = os.path.join(path, filename)
        
        total = 0
        
        self.__stream.seek(offset, 0)

        with open(filepath, "wb") as output:
            while (total + CHUNK) < size:
                data = self.__stream.read(CHUNK)
                
                if not data: break
                
                output.write(data)
                total += len(data)
                
            if total < size:
                data = self.__stream.read(size - total)
                output.write(data)

class PPTX(PowerPointFormat):

    @classmethod
    def is_valid_format(cls, filename):
        return zipfile.is_zipfile(filename)
    
    def _process(self, filename):

        self.__zipfile = zipfile.ZipFile(filename)
        
        n = 1

        for file in self.__zipfile.namelist():
            path, name = os.path.split(file)
            name, ext = os.path.splitext(name)
            
            if path == "ppt/media":
                filename = "{0}{1}{2}".format(self.basename, n, ext)
                
                self._files[filename] = file
                
                n += 1
                
    def _extract(self, name, newname="", path=CWD):

        filename = newname or name
        
        if not name in self._files:
            raise IOError("No such file")
        
        filepath = os.path.join(path, filename)
        
        total = 0
        
        file = self.__zipfile.open(self._files[name])
        
        with open(filepath, "wb") as output:
            while True:
                data = file.read(CHUNK)
                
                if not data: break
                
                output.write(data)
                total += len(data)


def PPTExtractor(filename):
    for cls in PowerPointFormat.__subclasses__():
        if cls.is_valid_format(filename):
            return cls(filename)
    raise InvalidFormat("{0} is not a PowerPoint file".format(filename))

def pdfimageextractor(directory, filename, pdf):
    i = 0
    njpg = 0
    while True:
        istream = pdf.find("stream", i)
        if istream < 0:
            break
        istart = pdf.find(startmark, istream, istream+20)
        if istart < 0:
            i = istream+20
            continue
        iend = pdf.find("endstream", istart)
        if iend < 0:
            raise Exception("Didn't find end of stream!")
        iend = pdf.find(endmark, iend-20)
        if iend < 0:
            raise Exception("Didn't find end of JPG!")
    
        istart += startfix
        iend += endfix

        jpg = pdf[istart:iend]
        jpgfile = file(directory+"\\"+filename+"-images\\jpg%d.jpg" % njpg, "wb")
        jpgfile.write(jpg)
        jpgfile.close()
    
        njpg += 1
        i = iend

def allimagestringextractor(imagesfoldername):
    allocrtext = ""
    for filename in os.listdir(imagesfoldername):
        if (filename.endswith(".jpg"))or(filename.endswith(".png")): 
            os.chdir(imagesfoldername)
            allocrtext = allocrtext + " " + pytesseract.image_to_string(Image.open(filename))
            continue
        else:
            continue
    return allocrtext

def convert_pdf_to_txt(path):
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
    fp = file(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    maxpages = 0
    caching = True
    pagenos=set()

    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password, caching=caching, check_extractable=True):
        interpreter.process_page(page)

    text = retstr.getvalue()

    fp.close()
    device.close()
    retstr.close()
    return text

def pdfstringextractor( pdffilename ):

   pdf2string = convert_pdf_to_txt(str(pdffilename))
   pdf2string = pdf2string.replace("\n", "")
   return pdf2string;

def pptstringextractor( pptfilename ):
   prs = Presentation(str(pptfilename))

   ppt2string = []

   for slide in prs.slides:
       for shape in slide.shapes:
           if not shape.has_text_frame:
               continue
           for paragraph in shape.text_frame.paragraphs:
               for run in paragraph.runs:
                   ppt2string.append(run.text)

   return str(ppt2string);

os.chdir('Drive:\\DataWareHouseFolderName\\')
directory = os.getcwd()

conn = psycopg2.connect("dbname=pythontest host=localhost user=<enter DB username> password=<enter DB password>")
cur = conn.cursor()

slideprocessingcount = 0

# Extract PDFs and PPTXs	
for filename in os.listdir(directory):
    if filename.endswith(".pdf"): 

        slideprocessingcount = slideprocessingcount + 1
        print("Processing file # " + str(slideprocessingcount) + " | TYPE: PDF")

        pdfcheck = "true"
        pptcheck = "false"

        completepdftext = pdfstringextractor(filename)
      		
        pdf = file(filename, "rb").read()
        os.mkdir(filename+'-images')
        pdfimageextractor(directory, filename, pdf)
        imagesfoldername = directory+'\\'+filename+'-images'
        allocrtext = allimagestringextractor(imagesfoldername)
        os.chdir(directory)
		
        cur.execute("INSERT INTO test (dwfilename, ispdf, isppt, completetext, allocrtext) VALUES (%s, %s, %s, %s, %s)",(str(filename), pdfcheck, pptcheck, completepdftext, str(allocrtext)))
        conn.commit()

        continue
    if filename.endswith(".pptx"):

        slideprocessingcount = slideprocessingcount + 1
        print("Processing file # " + str(slideprocessingcount) + " | TYPE: PPTX")

        pdfcheck = "false"
        pptcheck = "true"

        completeppttext = pptstringextractor(filename)

        ppt = PPTExtractor(filename)
        os.mkdir(filename+'-images')
        ppt.extractall(path=directory+'\\'+filename+'-images')
        imagesfoldername = directory+'\\'+filename+'-images'
        allocrtext = allimagestringextractor(imagesfoldername)
        os.chdir(directory)

        cur.execute("INSERT INTO test (dwfilename, ispdf, isppt, completetext, allocrtext) VALUES (%s, %s, %s, %s, %s)",(str(filename), pdfcheck, pptcheck, completeppttext, str(allocrtext)))
        conn.commit()

        continue
    else:
        continue


